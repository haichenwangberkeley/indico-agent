#!/usr/bin/env python3
from __future__ import annotations

import argparse
import datetime as dt
import hashlib
import json
import os
import re
import shutil
import sys
import urllib.error
import urllib.parse
import urllib.request
from dataclasses import dataclass
from pathlib import Path
from typing import Any


BASE_URL = "https://indico.cern.ch"
SIGNAL_TERMS = (
    "performance",
    "status",
    "recommendation",
    "efficiency",
    "scale factor",
    "resolution",
    "calibration",
    "run 3",
    "run 4",
    "improvement",
    "software",
    "reconstruction",
    "electron",
    "photon",
)
RESULT_TERMS = (
    "observed",
    "expected",
    "result",
    "limit",
    "significance",
    "measurement",
    "cross section",
    "confidence",
    "exclusion",
    "fit",
    "uncertainty",
    "systematic",
    "yield",
    "background",
    "signal",
    "paper",
    "conf note",
    "publication",
)
EDITORIAL_TERMS = (
    "editorial",
    "editor",
    "eb",
    "circulation",
    "approval",
    "approve",
    "comment",
    "question",
    "request",
    "recommendation",
    "concern",
    "change",
    "revision",
    "draft",
    "reader",
    "referee",
)
ACTION_TERMS = (
    "action",
    "todo",
    "to do",
    "follow up",
    "follow-up",
    "deadline",
    "by ",
    "assign",
    "responsible",
    "next step",
    "minutes",
    "decision",
    "conclusion",
)


@dataclass
class Material:
    title: str
    url: str
    talk_title: str = ""
    presenter: str = ""
    folder: str = ""
    content_type: str = ""
    event_id: str = ""
    event_title: str = ""
    event_date: str = ""
    checksum: str = ""
    modified_dt: str = ""
    size: int = 0


def request_bytes(url: str, token: str | None = None) -> tuple[bytes, dict[str, str], str]:
    headers = {"User-Agent": "cern-indico-briefing-agent/0.1"}
    if token:
        headers["Authorization"] = f"Bearer {token}"
    req = urllib.request.Request(url, headers=headers)
    try:
        with urllib.request.urlopen(req, timeout=120) as resp:
            return resp.read(), dict(resp.headers.items()), resp.geturl()
    except urllib.error.HTTPError as exc:
        body = exc.read().decode("utf-8", errors="replace")
        raise RuntimeError(f"HTTP {exc.code} for {url}: {body[:300]}") from exc


def request_json(url: str, token: str | None = None) -> Any:
    body, _, _ = request_bytes(url, token)
    return json.loads(body.decode("utf-8"))


def load_source_config(path: Path) -> dict[str, Any]:
    if not path.exists():
        raise RuntimeError(f"Source config not found: {path}")
    return json.loads(path.read_text(encoding="utf-8"))


def safe_name(name: str, default: str = "material") -> str:
    name = re.sub(r"[^\w.\-]+", "_", name.strip(), flags=re.ASCII).strip("._")
    return name or default


def full_url(url: str) -> str:
    return urllib.parse.urljoin(BASE_URL, url)


def people_names(entry: dict[str, Any]) -> str:
    names = []
    for person in (entry.get("presenters") or entry.get("speakers") or []):
        if isinstance(person, dict):
            names.append(person.get("name") or person.get("fullName") or person.get("full_name") or "")
    return ", ".join(n for n in names if n)


def collect_materials_from_entry(entry: dict[str, Any], event_info: dict[str, str] | None = None) -> list[Material]:
    materials: list[Material] = []
    event_info = event_info or {}
    talk_title = str(entry.get("title") or "")
    presenter = people_names(entry)

    def add_attachment(attachment: dict[str, Any], folder_title: str = "") -> None:
        url = attachment.get("download_url") or attachment.get("link_url")
        if not url:
            return
        materials.append(
            Material(
                title=str(attachment.get("title") or attachment.get("filename") or "attachment"),
                url=full_url(str(url)),
                talk_title=talk_title,
                presenter=presenter,
                folder=folder_title,
                content_type=str(attachment.get("content_type") or ""),
                event_id=event_info.get("event_id", ""),
                event_title=event_info.get("event_title", ""),
                event_date=event_info.get("event_date", ""),
                checksum=str(attachment.get("checksum") or ""),
                modified_dt=str(attachment.get("modified_dt") or ""),
                size=int(attachment.get("size") or 0),
            )
        )

    raw_attachments = entry.get("attachments") or {}
    folders = []
    direct_files = []
    if isinstance(raw_attachments, dict):
        folders.extend(raw_attachments.get("folders") or [])
        direct_files.extend(raw_attachments.get("files") or [])
    elif isinstance(raw_attachments, list):
        direct_files.extend(raw_attachments)

    folders.extend(entry.get("folders") or [])
    for attachment in direct_files:
        if isinstance(attachment, dict):
            add_attachment(attachment)

    for folder in folders:
        folder_title = str(folder.get("title") or "")
        for attachment in folder.get("attachments") or []:
            if not isinstance(attachment, dict):
                continue
            add_attachment(attachment, folder_title)
    return materials


def walk_timetable(node: Any, event_info: dict[str, str] | None = None) -> list[Material]:
    materials: list[Material] = []
    if isinstance(node, dict) and "_type" in node and str(node.get("_type")) in {"Conference", "Meeting", "Lecture"}:
        event_info = {
            "event_id": str(node.get("id") or ""),
            "event_title": str(node.get("title") or ""),
            "event_date": str((node.get("startDate") or {}).get("date") or ""),
        }
    if isinstance(node, dict):
        if "attachments" in node or "folders" in node:
            materials.extend(collect_materials_from_entry(node, event_info))
        for child in node.values():
            materials.extend(walk_timetable(child, event_info))
    elif isinstance(node, list):
        for child in node:
            materials.extend(walk_timetable(child, event_info))
    return materials


def fetch_event_materials(event_id: str, token: str | None) -> list[Material]:
    materials: list[Material] = []
    event_info: dict[str, str] = {"event_id": event_id}
    for url in (
        f"{BASE_URL}/export/event/{event_id}.json?detail=sessions&pretty=yes",
        f"{BASE_URL}/export/timetable/{event_id}.json?pretty=yes",
    ):
        data = request_json(url, token)
        results = data.get("results", {})
        if isinstance(results, list) and results:
            event = results[0]
            if isinstance(event, dict):
                event_info = {
                    "event_id": str(event.get("id") or event_id),
                    "event_title": str(event.get("title") or ""),
                    "event_date": str((event.get("startDate") or {}).get("date") or ""),
                }
        materials.extend(walk_timetable(results, event_info))

    deduped: list[Material] = []
    seen = set()
    for material in materials:
        if material.url in seen:
            continue
        seen.add(material.url)
        deduped.append(material)
    return deduped


def iso_date(days_from_today: int) -> str:
    return (dt.date.today() + dt.timedelta(days=days_from_today)).isoformat()


def fetch_category_events(category_id: str, token: str | None, date_from: str, date_to: str) -> list[dict[str, Any]]:
    params = urllib.parse.urlencode({"from": date_from, "to": date_to, "pretty": "yes"})
    data = request_json(f"{BASE_URL}/export/categ/{category_id}.json?{params}", token)
    events = data.get("results", [])
    return events if isinstance(events, list) else []


def materials_from_categories(
    category_ids: list[str],
    token: str | None,
    date_from: str,
    date_to: str,
) -> list[Material]:
    materials: list[Material] = []
    seen_events = set()
    for category_id in category_ids:
        for event in fetch_category_events(category_id, token, date_from, date_to):
            event_id = str(event.get("id") or "")
            if not event_id or event_id in seen_events:
                continue
            seen_events.add(event_id)
            materials.extend(fetch_event_materials(event_id, token))
    return materials


def extract_pdf(path: Path) -> str:
    try:
        import pdfplumber
    except ImportError as exc:
        raise RuntimeError("pdfplumber is required for PDF extraction") from exc
    chunks = []
    with pdfplumber.open(path) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                chunks.append(f"\n\n--- Page {page_num} ---\n{text.strip()}")
    return "".join(chunks).strip()


def extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PPTX extraction") from exc
    prs = Presentation(path)
    chunks = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
        if lines:
            chunks.append(f"\n\n--- Slide {idx} ---\n" + "\n".join(lines))
    return "".join(chunks).strip()


def extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    return ""


def looks_like_pdf(path: Path) -> bool:
    with path.open("rb") as fh:
        return fh.read(5) == b"%PDF-"


def material_filename(material: Material) -> str:
    filename = safe_name(Path(urllib.parse.urlparse(material.url).path).name or material.title)
    if "." not in filename and material.title:
        filename = safe_name(material.title)
    if "." not in filename:
        if "pdf" in material.content_type.lower():
            filename = f"{filename}.pdf"
        elif "presentation" in material.content_type.lower() or "powerpoint" in material.content_type.lower():
            filename = f"{filename}.pptx"
    return filename


def cache_key(material: Material) -> str:
    stable = material.checksum or material.url
    return hashlib.sha256(stable.encode("utf-8")).hexdigest()[:24]


def download_material(material: Material, out_dir: Path, token: str | None, cache_dir: Path | None = None) -> Path:
    filename = material_filename(material)
    path = out_dir / filename
    if path.exists():
        path = out_dir / f"{path.stem}_{safe_name(material.event_id or 'copy')}{path.suffix}"

    cached_path = None
    if cache_dir:
        cache_dir.mkdir(parents=True, exist_ok=True)
        cached_path = cache_dir / f"{cache_key(material)}_{filename}"
        if cached_path.exists():
            shutil.copy2(cached_path, path)
            return path

    body, headers, final_url = request_bytes(material.url, token)
    path.write_bytes(body)
    content_type = headers.get("Content-Type", "")
    if "text/html" in content_type.lower():
        raise RuntimeError(f"Downloaded HTML instead of material from {final_url}")
    if cached_path:
        shutil.copy2(path, cached_path)
    return path


def is_supported_material(material: Material) -> bool:
    haystack = " ".join((material.title, material.url, material.content_type, material.folder)).lower()
    return any(x in haystack for x in (".pdf", ".pptx", ".ppt", "application/pdf", "powerpoint", "presentation"))


def signal_lines(text: str, limit: int = 16) -> list[str]:
    return matched_lines(text, SIGNAL_TERMS, limit)


def matched_lines(text: str, terms: tuple[str, ...], limit: int = 12) -> list[str]:
    lines = []
    seen = set()
    for raw in text.splitlines():
        line = re.sub(r"\s+", " ", raw).strip()
        if len(line) < 20 or line in seen:
            continue
        lower = line.lower()
        if any(term in lower for term in terms):
            lines.append(line)
            seen.add(line)
        if len(lines) >= limit:
            break
    return lines


def prose_from_matches(matches: list[str]) -> str:
    if not matches:
        return "No keyword-matched lines were found in the extracted text."
    sentences = []
    for item in matches:
        item = item.rstrip(".")
        sentences.append(f"{item}.")
    return " ".join(sentences)


def section(lines: list[str], title: str, matches: list[str]) -> None:
    lines.append(title)
    lines.append("")
    lines.append(prose_from_matches(matches))
    lines.append("")


def write_empty_briefing(run_dir: Path, reason: str) -> Path:
    now = dt.datetime.now(dt.timezone.utc).astimezone().isoformat(timespec="seconds")
    lines = [
        "# CERN Indico Briefing",
        "",
        f"Generated: {now}",
        "",
        "No Processable Materials",
        "",
        reason,
        "",
    ]
    path = run_dir / "briefing.md"
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def write_briefing(run_dir: Path, processed: list[tuple[Material, Path, Path | None, str]]) -> Path:
    now = dt.datetime.now(dt.timezone.utc).astimezone().isoformat(timespec="seconds")
    lines = [
        "# CERN Indico Briefing",
        "",
        f"Generated: {now}",
        "",
        "## Materials",
        "",
    ]
    for idx, (material, file_path, text_path, text) in enumerate(processed, start=1):
        lines.append(f"### {idx}. {material.talk_title or material.title}")
        context = []
        if material.event_title:
            event = material.event_title
            if material.event_date:
                event = f"{event} on {material.event_date}"
            context.append(f"It comes from {event}")
        if material.presenter:
            context.append(f"presented by {material.presenter}")
        if material.folder:
            context.append(f"from the {material.folder} material folder")
        if context:
            lines.append(" ".join(context) + ".")
            lines.append("")
        artifact_sentence = f"The source is {material.url}. The local file is `{file_path.relative_to(run_dir)}`."
        if text_path:
            artifact_sentence += f" The extracted text is `{text_path.relative_to(run_dir)}`."
        lines.append(artifact_sentence)
        lines.append("")
        section(lines, "Main Results:", matched_lines(text, RESULT_TERMS, 12))
        section(lines, "Editorial Board Reactions:", matched_lines(text, EDITORIAL_TERMS, 10))
        section(lines, "Minutes And Action Items:", matched_lines(text, ACTION_TERMS, 10))
        section(lines, "Other Signals:", signal_lines(text, 8))
        lines.append("")
    path = run_dir / "briefing.md"
    path.write_text("\n".join(lines), encoding="utf-8")
    return path


def process_materials(
    materials: list[Material],
    out_root: Path,
    token: str | None,
    match: str = "",
    limit: int | None = None,
) -> Path:
    stamp = dt.datetime.now().strftime("%Y%m%d_%H%M%S")
    run_dir = out_root / stamp
    material_dir = run_dir / "materials"
    text_dir = run_dir / "text"
    cache_dir = out_root.parent / "cache" / "materials"
    material_dir.mkdir(parents=True, exist_ok=True)
    text_dir.mkdir(parents=True, exist_ok=True)

    if match:
        needle = match.lower()
        materials = [
            material
            for material in materials
            if needle in " ".join((material.title, material.talk_title, material.presenter, material.url)).lower()
        ]

    processed: list[tuple[Material, Path, Path | None, str]] = []
    for material in materials:
        if limit is not None and len(processed) >= limit:
            break
        if not is_supported_material(material):
            continue
        file_path = download_material(material, material_dir, token, cache_dir)
        text = ""
        text_path = None
        if file_path.suffix.lower() in {".pdf", ".pptx"}:
            text = extract_text(file_path)
            text_path = text_dir / f"{file_path.stem}.txt"
            text_path.write_text(text, encoding="utf-8")
        processed.append((material, file_path, text_path, text))

    if not processed:
        return write_empty_briefing(
            run_dir,
            "The configured Indico sources were checked, but no PDF, PPTX, or PPT materials matched the current window and filters.",
        )
    return write_briefing(run_dir, processed)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Fetch CERN Indico materials and write a text briefing.")
    parser.add_argument("--event", help="Indico event ID to process from timetable export.")
    parser.add_argument("--category", action="append", help="Indico category ID to discover events from. May be passed more than once.")
    parser.add_argument("--source-config", help="JSON source config with category_ids and defaults.")
    parser.add_argument("--attachment-url", help="Single Indico attachment URL to process.")
    parser.add_argument("--title", default="", help="Title to use with --attachment-url.")
    parser.add_argument("--out", default="output/briefings", help="Output directory.")
    parser.add_argument("--match", default="", help="Only process materials whose title, talk, presenter, or URL contains this text.")
    parser.add_argument("--limit", type=int, default=20, help="Maximum number of matching materials to process.")
    parser.add_argument("--from-date", help="Category discovery start date, YYYY-MM-DD.")
    parser.add_argument("--to-date", help="Category discovery end date, YYYY-MM-DD.")
    parser.add_argument("--lookback-days", type=int, default=7, help="Category discovery lookback window.")
    parser.add_argument("--lookahead-days", type=int, default=1, help="Category discovery lookahead window.")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    token = os.environ.get("INDICO_TOKEN")
    categories = list(args.category or [])
    if args.source_config:
        config = load_source_config(Path(args.source_config))
        categories.extend(str(x) for x in config.get("category_ids", []))
        args.match = args.match or str(config.get("match", ""))
        args.limit = args.limit if args.limit is not None else config.get("limit")
        args.lookback_days = int(config.get("lookback_days", args.lookback_days))
        args.lookahead_days = int(config.get("lookahead_days", args.lookahead_days))

    if not args.event and not args.attachment_url and not categories:
        print("Pass --event, --category, --source-config, or --attachment-url", file=sys.stderr)
        return 2

    if args.attachment_url:
        materials = [Material(title=args.title or Path(urllib.parse.urlparse(args.attachment_url).path).name, url=args.attachment_url)]
    elif args.event:
        materials = fetch_event_materials(args.event, token)
    else:
        date_from = args.from_date or iso_date(-args.lookback_days)
        date_to = args.to_date or iso_date(args.lookahead_days)
        materials = materials_from_categories(categories, token, date_from, date_to)

    briefing = process_materials(materials, Path(args.out), token, match=args.match, limit=args.limit)
    print(briefing)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
